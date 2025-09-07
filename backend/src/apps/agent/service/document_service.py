"""
文档处理服务 - 负责文件解析和内容提取
"""
import os
import uuid
import json
import csv
import io
import subprocess
import tempfile
from typing import Dict, Any, List, Optional
from datetime import datetime
from pathlib import Path
from sqlalchemy import select, update
from sqlalchemy.ext.asyncio import AsyncSession

# 文档解析库
try:
    from pypdf import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

from src.shared.core.logging import get_logger

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logger = get_logger(__name__)

# 在模块加载时检查依赖并给出警告
if not HAS_PDF:
    logger.warning("未安装 pypdf 库，PDF 文件解析功能将不可用。请运行: pip install pypdf")
    
if not HAS_DOCX:
    logger.warning("未安装 python-docx 库，Word 文档解析功能将不可用。请运行: pip install python-docx")
    
if not HAS_OPENPYXL:
    logger.warning("未安装 openpyxl 库，Excel 文件解析功能将不可用。请运行: pip install openpyxl")
    
if not HAS_PPTX:
    logger.warning("未安装 python-pptx 库，PowerPoint 文件解析功能将不可用。请运行: pip install python-pptx")
from src.shared.core.exceptions import BusinessException
from src.shared.schemas.response import ResponseCode
from src.shared.db.models import now_shanghai
from src.shared.core.config import settings
from ..models import AgentDocumentUpload

# 从配置中获取文件存储路径
UPLOAD_DIR = Path(settings.UPLOAD_DIR)
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

# 创建临时目录（用于存储转换的临时文件）
TEMP_DIR = UPLOAD_DIR.parent / 'temp'
TEMP_DIR.mkdir(parents=True, exist_ok=True)
logger.info(f"临时文件目录: {TEMP_DIR}")

# 从配置中获取支持的文件类型
SUPPORTED_FILE_TYPES = settings.UPLOAD_ALLOWED_EXTENSIONS

# 从配置中获取文件大小限制
MAX_FILE_SIZE = settings.MAX_UPLOAD_SIZE_MB * 1024 * 1024

# 处理状态枚举
class ProcessStatus:
    UPLOADED = 0
    PROCESSING = 1
    READY = 2
    FAILED = 3


class DocumentService:
    """文档处理服务"""
    
    def __init__(self):
        """初始化文档处理器映射"""
        # 文件处理器映射表 - 策略模式
        self._file_processors = {
            # 文本文件
            '.txt': self._process_text_file,
            '.md': self._process_text_file,
            
            # PDF文件
            '.pdf': self._process_pdf_file if HAS_PDF else None,
            
            # Word文档
            '.doc': self._process_doc_file,
            '.docx': self._process_docx_file if HAS_DOCX else None,
            
            # PowerPoint演示文稿
            '.pptx': self._process_pptx_file if HAS_PPTX else None,
            
            # 表格文件
            '.csv': self._process_csv_file,
            '.xlsx': self._process_excel_file if HAS_OPENPYXL else None,
            '.xls': self._process_excel_file if HAS_OPENPYXL else None,
            
            # 图片文件
            '.png': self._process_image_file,
            '.jpg': self._process_image_file,
            '.jpeg': self._process_image_file,
            '.gif': self._process_image_file,
            '.bmp': self._process_image_file,
            '.webp': self._process_image_file,
            
            # 代码和配置文件 - 作为文本处理
            '.sql': self._process_text_file,
            '.yaml': self._process_text_file,
            '.yml': self._process_text_file,
            '.json': self._process_text_file,
            '.xml': self._process_text_file,
            '.ini': self._process_text_file,
            '.conf': self._process_text_file,
            '.cfg': self._process_text_file,
            '.properties': self._process_text_file,
            '.env': self._process_text_file,
            '.toml': self._process_text_file,
            
            # 编程语言文件 - 作为文本处理
            '.py': self._process_text_file,
            '.js': self._process_text_file,
            '.ts': self._process_text_file,
            '.jsx': self._process_text_file,
            '.tsx': self._process_text_file,
            '.java': self._process_text_file,
            '.cpp': self._process_text_file,
            '.c': self._process_text_file,
            '.h': self._process_text_file,
            '.cs': self._process_text_file,
            '.php': self._process_text_file,
            '.rb': self._process_text_file,
            '.go': self._process_text_file,
            '.rs': self._process_text_file,
            '.swift': self._process_text_file,
            '.kt': self._process_text_file,
            
            # 脚本文件 - 作为文本处理
            '.sh': self._process_text_file,
            '.bash': self._process_text_file,
            '.zsh': self._process_text_file,
            '.ps1': self._process_text_file,
            '.bat': self._process_text_file,
            '.cmd': self._process_text_file,
            
            # Web相关文件 - 作为文本处理
            '.html': self._process_text_file,
            '.htm': self._process_text_file,
            '.css': self._process_text_file,
            '.scss': self._process_text_file,
            '.sass': self._process_text_file,
            '.less': self._process_text_file,
            
            # 日志和文档文件 - 作为文本处理
            '.log': self._process_text_file,
            '.out': self._process_text_file,
            '.err': self._process_text_file,
            '.trace': self._process_text_file,
            '.rst': self._process_text_file,
            '.org': self._process_text_file,
            '.tex': self._process_text_file,
            '.rtf': self._process_text_file,
        }
    
    async def upload_file(self, db: AsyncSession, file_content: bytes, filename: str, user_name: str) -> Dict[str, Any]:
        """
        上传文件
        
        Args:
            db: 数据库会话
            file_content: 文件内容
            filename: 文件名
            user_name: 用户名
            
        Returns:
            文件信息
        """
        # 检查文件大小
        if len(file_content) > MAX_FILE_SIZE:
            raise BusinessException(
                f"文件大小超过限制（最大{settings.MAX_UPLOAD_SIZE_MB}MB）",
                ResponseCode.BAD_REQUEST
            )
        
        # 检查文件类型
        file_ext = Path(filename).suffix.lower()
        if file_ext not in SUPPORTED_FILE_TYPES:
            raise BusinessException(
                f"不支持的文件类型：{file_ext}",
                ResponseCode.BAD_REQUEST
            )
        
        # 生成文件ID和保存路径
        file_id = str(uuid.uuid4())
        # 安全处理：只使用生成的文件ID和扩展名，避免路径遍历
        safe_filename = f"{file_id}{file_ext}"
        file_path = UPLOAD_DIR / safe_filename
        
        # 保存文件
        import aiofiles
        async with aiofiles.open(file_path, 'wb') as f:
            await f.write(file_content)
        
        # 保存到数据库
        async with db.begin():
            doc_upload = AgentDocumentUpload(
                file_id=file_id,
                file_name=filename,
                file_size=len(file_content),
                file_type=file_ext,
                file_path=str(file_path),
                process_status=ProcessStatus.UPLOADED,
                create_by=user_name,
                update_by=user_name
            )
            db.add(doc_upload)
            await db.flush()
            
            # 获取生成的ID
            doc_id = doc_upload.id
        
        logger.info(f"文件上传成功: {filename} -> {file_id}")
        
        # 异步处理文档 - 使用后台任务，避免事务冲突
        import asyncio
        # 创建任务并保存引用，避免任务丢失
        task = asyncio.create_task(self._process_document_async(file_id))
        # 添加错误处理回调
        task.add_done_callback(lambda t: logger.error(f"文档处理任务异常: {t.exception()}") if t.exception() else None)
        
        return {
            "file_id": file_id,
            "file_name": filename,
            "file_size": len(file_content),
            "file_type": file_ext,
            "upload_time": now_shanghai().isoformat(),
            "status": "uploaded"
        }
    
    async def process_document(self, db: AsyncSession, file_id: str) -> None:
        """
        处理文档，提取内容
        
        Args:
            db: 数据库会话
            file_id: 文件ID
        """
        try:
            async with db.begin():
                # 查询文档记录
                result = await db.execute(
                    select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
                )
                doc_upload = result.scalar_one_or_none()
                
                if not doc_upload:
                    raise BusinessException("文件不存在", ResponseCode.NOT_FOUND)
                
                # 更新状态为处理中
                doc_upload.process_status = ProcessStatus.PROCESSING
                doc_upload.process_start_time = now_shanghai()
                await db.flush()
            
            file_path = Path(doc_upload.file_path)
            file_ext = doc_upload.file_type
            file_name = doc_upload.file_name
            
            # 根据文件类型提取文本 - 使用策略模式
            full_text = await self._dispatch_file_processing(file_path, file_ext, file_name, file_id)
            
            # 如果没有提取到内容
            if not full_text.strip():
                full_text = f"[{file_name}]\n\n文档内容为空或无法提取文本内容。"
            
            # 处理文档内容，分块并更新数据库
            char_count = await self._finalize_document_processing(db, file_id, full_text, file_ext)
            
            logger.info(f"文档处理完成: {file_id}, 类型: {file_ext}, 字符数: {char_count}")
            
        except Exception as e:
            logger.error(f"文档处理失败: {file_id}, 错误: {str(e)}")
            # 更新失败状态 - 使用新的数据库会话
            from src.shared.db.config import get_async_db_context
            async with get_async_db_context() as error_db:
                result = await error_db.execute(
                    select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
                )
                doc_upload = result.scalar_one_or_none()
                if doc_upload:
                    doc_upload.process_status = ProcessStatus.FAILED
                    doc_upload.error_message = str(e)
                    doc_upload.process_end_time = now_shanghai()
                    await error_db.commit()
    
    async def get_document_content(self, db: AsyncSession, file_id: str, user_name: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """
        获取文档内容
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            user_name: 用户名（用于权限检查）
            
        Returns:
            文档内容
        """
        # 查询文档
        query = select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
        
        # 如果提供了用户名，检查所有权
        if user_name:
            query = query.where(AgentDocumentUpload.create_by == user_name)
        
        result = await db.execute(query)
        doc_upload = result.scalar_one_or_none()
        
        if not doc_upload:
            return None
        
        if doc_upload.process_status != ProcessStatus.READY:
            return None
        
        # 解析JSON字段
        chunks = []
        metadata = {}
        try:
            if doc_upload.doc_chunks:
                chunks = json.loads(doc_upload.doc_chunks)
            if doc_upload.doc_metadata:
                metadata = json.loads(doc_upload.doc_metadata)
        except Exception as e:
            logger.error(f"解析JSON字段失败: {e}")
        
        return {
            "file_id": file_id,
            "file_name": doc_upload.file_name,
            "file_size": doc_upload.file_size,
            "content": doc_upload.doc_content or "",
            "metadata": metadata,
            "chunks": chunks
        }
    
    async def get_file_status(self, db: AsyncSession, file_id: str) -> Optional[Dict[str, Any]]:
        """
        获取文件处理状态
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            
        Returns:
            文件状态
        """
        result = await db.execute(
            select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
        )
        doc_upload = result.scalar_one_or_none()
        
        if not doc_upload:
            return None
        
        # 状态映射
        status_map = {
            ProcessStatus.UPLOADED: "uploaded",
            ProcessStatus.PROCESSING: "processing",
            ProcessStatus.READY: "ready",
            ProcessStatus.FAILED: "failed"
        }
        
        return {
            "file_id": file_id,
            "status": status_map.get(doc_upload.process_status, "unknown"),
            "message": doc_upload.error_message,
            "processed_at": doc_upload.process_end_time.isoformat() if doc_upload.process_end_time else None
        }
    
    def get_documents_info(self, file_ids: List[str]) -> List[Dict[str, Any]]:
        """
        获取文档的元信息（不包含内容）
        
        Args:
            file_ids: 文件ID列表
            
        Returns:
            文档信息列表，每个元素包含 file_id, file_name, file_size 等
        """
        from src.shared.db.config import get_sync_db
        from ..models import AgentDocumentUpload
        
        docs_info = []
        
        db_gen = get_sync_db()
        db = next(db_gen)
        try:
            for file_id in file_ids:
                try:
                    # 查询文档
                    result = db.execute(
                        select(AgentDocumentUpload).where(
                            AgentDocumentUpload.file_id == file_id
                        )
                    )
                    document = result.scalar_one_or_none()
                    
                    if document and document.process_status == 2:  # READY
                        doc_info = {
                            "file_id": document.file_id,
                            "file_name": document.file_name,
                            "file_size": document.file_size,
                            "file_type": document.file_type,
                            "upload_time": str(document.upload_time)
                        }
                        docs_info.append(doc_info)
                        logger.debug(f"获取文档元信息: {doc_info['file_name']}")
                    else:
                        logger.warning(f"文档 {file_id} 不存在或未就绪")
                        
                except Exception as e:
                    logger.error(f"获取文档 {file_id} 元信息失败: {e}")
                    
            return docs_info
            
        finally:
            try:
                next(db_gen)
            except StopIteration:
                pass

    async def get_documents_info_async(self, db: AsyncSession, file_ids: List[str]) -> List[Dict[str, Any]]:
        """
        异步获取文档的元信息（不包含内容）
        
        Args:
            db: 异步数据库会话
            file_ids: 文件ID列表
            
        Returns:
            文档信息列表，每个元素包含 file_id, file_name, file_size 等
        """
        from ..models import AgentDocumentUpload
        
        docs_info = []
        
        for file_id in file_ids:
            try:
                # 异步查询文档
                result = await db.execute(
                    select(AgentDocumentUpload).where(
                        AgentDocumentUpload.file_id == file_id
                    )
                )
                document = result.scalar_one_or_none()
                
                if document and document.process_status == 2:  # READY
                    doc_info = {
                        "file_id": document.file_id,
                        "file_name": document.file_name,
                        "file_size": document.file_size,
                        "file_type": document.file_type,
                        "upload_time": str(document.upload_time)
                    }
                    docs_info.append(doc_info)
                    logger.debug(f"✅ 获取文档元信息成功: {doc_info['file_name']} (ID: {file_id})")
                elif document and document.process_status != 2:
                    logger.warning(f"⚠️ 文档 {file_id} 尚未就绪，状态: {document.process_status}")
                else:
                    logger.warning(f"❌ 文档 {file_id} 不存在")
                    
            except Exception as e:
                logger.error(f"获取文档 {file_id} 元信息失败: {e}")
                
        return docs_info

    def get_document_context(self, file_ids: List[str], max_length: int = 10000000) -> str:
        """
        获取文档上下文（用于对话）
        注意：这是同步方法，因为在 streaming.py 中被同步调用
        
        Args:
            file_ids: 文件ID列表
            max_length: 最大长度限制，太小会导致前端看不到对应的文档,这里不做控制，在模型处控制内容大小吧
            
        Returns:
            文档上下文文本
        """
        from src.shared.db.config import get_sync_db
        
        context_parts = []
        current_length = 0
        
        db_gen = get_sync_db()
        db = next(db_gen)
        try:
            for file_id in file_ids:
                result = db.query(AgentDocumentUpload).filter(
                    AgentDocumentUpload.file_id == file_id,
                    AgentDocumentUpload.process_status == ProcessStatus.READY
                ).first()
                
                if not result:
                    continue
                
                content = result.doc_content or ""
                file_name = result.file_name
                
                # 如果内容太长，只取前面部分
                if current_length + len(content) > max_length:
                    remaining = max_length - current_length
                    if remaining > 100:  # 至少保留100字符
                        content = content[:remaining] + "内容过长，只获取部分信息，要通知用户，部分信息已丢失"
                    else:
                        break
                
                # 为每个文档添加明确的开始和结束标记
                doc_with_markers = f"""
==================== 文档开始 ====================
文档名称：{file_name}
文档ID：{file_id}
--------------------------------------------------
{content}
==================== 文档结束 ===================="""
                context_parts.append(doc_with_markers.strip())
                current_length += len(content)
        finally:
            db.close()
        
        # 如果有多个文档，在整体内容前添加文档摘要
        if len(context_parts) > 1:
            header = f"共有 {len(context_parts)} 个文档供参考。每个文档都用明确的分隔符标记了开始和结束。\n"
            return header + "\n\n".join(context_parts)
        elif len(context_parts) == 1:
            return context_parts[0]
        else:
            return ""


    async def _process_text_file(self, file_path: Path, file_name: str) -> str:
        """处理文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"文本文件读取错误: {e}")
            return f"文本文件读取失败: {str(e)}"
    
    async def _process_pdf_file(self, file_path: Path, file_name: str) -> str:
        """处理PDF文件"""
        try:
            reader = PdfReader(str(file_path))
            text_parts = []
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    text_parts.append(f"[第{page_num + 1}页]\n{text}")
            full_text = "\n\n".join(text_parts)
            logger.info(f"PDF解析成功: {file_name}, 页数: {len(reader.pages)}")
            return full_text
        except Exception as e:
            logger.error(f"PDF解析错误: {e}")
            return f"PDF文件解析失败: {str(e)}"
    
    async def _process_doc_file(self, file_path: Path, file_name: str, file_id: str) -> tuple[str, Path, str]:
        """处理.doc文件，返回处理结果和可能更新的文件路径"""
        logger.info(f"开始转换 .doc 文件: {file_name}")
        try:
            # 使用固定的临时目录，避免自动清理
            # 为避免文件名冲突，使用时间戳和文件ID
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            temp_subdir = TEMP_DIR / f"doc_conversion_{timestamp}_{file_id[:8]}"
            temp_subdir.mkdir(parents=True, exist_ok=True)
            
            # 使用 LibreOffice 转换
            cmd = ['soffice', '--headless', '--convert-to', 'docx', '--outdir', str(temp_subdir), str(file_path)]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                # 转换成功，更新文件路径为临时的 docx 文件
                temp_docx = temp_subdir / (file_path.stem + '.docx')
                if temp_docx.exists():
                    logger.info(f".doc 转换成功，临时文件位置: {temp_docx}")
                    logger.info("提示：临时文件保留在 temp 目录中，可以手动清理")
                    return "", temp_docx, '.docx'
                else:
                    raise Exception("转换后的文件未找到")
            else:
                raise Exception(f"LibreOffice 转换失败: {result.stderr}")
        except FileNotFoundError:
            # LibreOffice未安装的情况
            logger.error(f".doc 转换失败: LibreOffice未安装")
            return f"[.doc 文件: {file_name}]\n\n文档解析失败：系统未安装 LibreOffice，无法处理 .doc 格式文件。请联系管理员安装 LibreOffice 或将文件转换为 .docx 格式后重新上传。\n\n请通知管理员查看系统日志以了解详细原因。", file_path, '.doc'
        except subprocess.TimeoutExpired:
            # 转换超时的情况
            logger.error(f".doc 转换失败: 转换超时")
            return f"[.doc 文件: {file_name}]\n\n文档解析失败：.doc 文件转换超时，可能文件过大或系统繁忙。请尝试将文件转换为 .docx 格式后重新上传。\n\n请通知管理员查看系统日志以了解详细原因。", file_path, '.doc'
        except Exception as e:
            # 其他错误情况
            logger.error(f".doc 转换失败: {e}")
            error_msg = str(e).lower()
            if 'no such file' in error_msg or 'not found' in error_msg:
                return f"[.doc 文件: {file_name}]\n\n文档解析失败：系统未安装 LibreOffice，无法处理 .doc 格式文件。请联系管理员安装 LibreOffice 或将文件转换为 .docx 格式后重新上传。\n\n请通知管理员查看系统日志以了解详细原因。", file_path, '.doc'
            else:
                return f"[.doc 文件: {file_name}]\n\n文档解析失败：.doc 文件转换出错 - {str(e)}。建议将文件转换为 .docx 格式后重新上传。\n\n请通知管理员查看系统日志以了解详细原因。", file_path, '.doc'
    
    async def _process_docx_file(self, file_path: Path, file_name: str) -> str:
        """处理Word文档"""
        try:
            doc = docx.Document(str(file_path))
            content_parts = []
            
            # 导入必要的类型用于按顺序遍历文档元素
            from docx.document import Document as _Document
            from docx.oxml.text.paragraph import CT_P
            from docx.oxml.table import CT_Tbl
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
            
            def iter_block_items(parent):
                """按文档顺序遍历段落和表格"""
                if isinstance(parent, _Document):
                    parent_elm = parent.element.body
                elif isinstance(parent, _Cell):
                    parent_elm = parent._tc
                else:
                    raise ValueError("something's not right")
                
                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        yield Paragraph(child, parent)
                    elif isinstance(child, CT_Tbl):
                        yield Table(child, parent)
            
            # 第一阶段：按文档顺序遍历，收集内容和图片位置
            images_to_process = []  # 待处理的图片列表
            image_placeholders = {}  # 占位符映射
            image_count = 0
            
            # 首先收集所有图片的数据
            image_data_by_id = {}
            for rel_id, rel in doc.part.rels.items():
                if "image" in rel.reltype:
                    image_data_by_id[rel_id] = rel.target_part.blob
            
            # 按文档顺序遍历元素
            for block in iter_block_items(doc):
                if isinstance(block, Paragraph):
                    # 检查段落中是否有图片
                    has_image = False
                    para_text = block.text.strip()
                    
                    # 简化的图片检测方法 - 检查段落的_element是否包含drawing
                    try:
                        # 获取段落的XML字符串表示
                        para_xml = block._element.xml
                        if 'w:drawing' in para_xml:
                            # 这个段落包含图片，查找所有的r:embed属性
                            import re
                            # 使用正则表达式查找所有的r:embed引用
                            embed_pattern = r'r:embed="(rId\d+)"'
                            embed_matches = re.findall(embed_pattern, para_xml)
                            
                            for rel_id in embed_matches:
                                if rel_id in image_data_by_id:
                                    has_image = True
                                    image_count += 1
                                    placeholder = f"[[IMAGE_PLACEHOLDER_{image_count}]]"
                                    
                                    # 记录图片信息
                                    images_to_process.append({
                                        'index': image_count,
                                        'placeholder': placeholder,
                                        'data': image_data_by_id[rel_id]
                                    })
                                    
                                    # 如果段落有文字，先添加文字
                                    if para_text:
                                        content_parts.append(para_text)
                                        para_text = ""  # 清空以避免重复添加
                                    
                                    # 添加图片占位符
                                    content_parts.append(placeholder)
                    except Exception as e:
                        logger.debug(f"检查段落图片时出错: {e}")
                    
                    # 如果段落没有图片但有文字，添加文字
                    if not has_image and para_text:
                        content_parts.append(para_text)
                        
                elif isinstance(block, Table):
                    # 处理表格
                    table_text = "\n[表格]\n"
                    for row in block.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        table_text += row_text + "\n"
                    content_parts.append(table_text)
            
            # 第二阶段：如果有图片需要处理，并发处理所有图片
            if images_to_process and settings.VISION_API_KEY:
                from .multimodal_service import multimodal_service
                logger.info(f"Word文档包含 {len(images_to_process)} 张图片，开始并发处理...")
                
                # 创建并发任务
                async def process_single_image(image_info):
                    """处理单张图片的异步函数"""
                    index = image_info['index']
                    placeholder = image_info['placeholder']
                    try:
                        result = await multimodal_service.extract_image_content(
                            image_info['data'],
                            options={
                                'prompt': '请详细描述这张图片的内容，如果是图表请提取其中的数据和信息。'
                            }
                        )
                        
                        if result['success']:
                            logger.info(f"成功提取图片 {index} 内容")
                            return {
                                'placeholder': placeholder,
                                'content': f"\n[图片 {index}]\n{result['content']}\n",
                                'success': True
                            }
                        else:
                            logger.warning(f"图片 {index} 提取失败: {result.get('error', '未知错误')}")
                            return {
                                'placeholder': placeholder,
                                'content': f"\n[图片 {index}：无法识别 - {result.get('error', '未知错误')}]\n",
                                'success': False
                            }
                    except Exception as e:
                        logger.error(f"处理图片 {index} 时出错: {e}")
                        return {
                            'placeholder': placeholder,
                            'content': f"\n[图片 {index}：处理失败 - {str(e)}]\n",
                            'success': False
                        }
                
                # 限制并发数，避免API限流
                import asyncio
                MAX_CONCURRENT = settings.IMAGE_PROCESS_MAX_CONCURRENT
                semaphore = asyncio.Semaphore(MAX_CONCURRENT)
                
                async def process_with_semaphore(image_info):
                    async with semaphore:
                        try:
                            return await asyncio.wait_for(
                                process_single_image(image_info), 
                                timeout=settings.IMAGE_PROCESS_TIMEOUT
                            )
                        except asyncio.TimeoutError:
                            logger.error(f"处理图片 {image_info['index']} 超时")
                            return {
                                'placeholder': image_info['placeholder'],
                                'content': f"\n[图片 {image_info['index']}：处理超时]\n",
                                'success': False
                            }
                
                # 并发处理所有图片
                tasks = [process_with_semaphore(img) for img in images_to_process]
                results = await asyncio.gather(*tasks, return_exceptions=True)
                
                # 建立占位符到内容的映射
                placeholder_map = {}
                success_count = 0
                for result in results:
                    if isinstance(result, Exception):
                        logger.error(f"图片处理任务异常: {result}")
                    elif isinstance(result, dict):
                        placeholder_map[result['placeholder']] = result['content']
                        if result.get('success'):
                            success_count += 1
                
                logger.info(f"图片并发处理完成，成功处理 {success_count} 张图片")
                
                # 第三阶段：替换占位符为实际内容
                for i, part in enumerate(content_parts):
                    if part in placeholder_map:
                        content_parts[i] = placeholder_map[part]
            elif images_to_process and not settings.VISION_API_KEY:
                # 如果有图片但没有配置视觉模型，替换占位符为提示信息
                for i, part in enumerate(content_parts):
                    if part.startswith("[[IMAGE_PLACEHOLDER_"):
                        image_num = part.split('_')[2].rstrip(']]')
                        content_parts[i] = f"\n[图片 {image_num}：未配置视觉模型]\n"
            
            full_text = "\n\n".join(content_parts)
            logger.info(f"Word文档解析成功: {file_name}")
            return full_text
            
        except Exception as e:
            logger.error(f"Word文档解析错误: {e}")
            return f"Word文档解析失败: {str(e)}"
    
    async def _process_csv_file(self, file_path: Path, file_name: str) -> str:
        """处理CSV文件"""
        try:
            content_parts = ["[CSV文件]"]
            
            # 尝试不同的编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'cp1252', 'iso-8859-1']
            csv_content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        csv_content = f.read()
                        used_encoding = encoding
                        break
                except UnicodeDecodeError:
                    continue
            
            if csv_content is None:
                raise Exception("无法识别文件编码")
            
            # 使用csv模块解析
            csv_reader = csv.reader(io.StringIO(csv_content))
            rows = list(csv_reader)
            
            if rows:
                # 添加编码信息
                content_parts.append(f"编码: {used_encoding}")
                content_parts.append(f"行数: {len(rows)}")
                content_parts.append("")
                
                # 将CSV内容转换为表格格式
                for i, row in enumerate(rows):
                    if i == 0:  # 如果第一行是表头，可以特殊标记
                        content_parts.append("[表头]")
                    row_text = " | ".join([cell.strip() for cell in row])
                    content_parts.append(row_text)
                    
                    # 对于大文件，限制显示行数
                    if i > 1000 and len(rows) > 1100:
                        content_parts.append(f"\n... 省略 {len(rows) - 1001} 行 ...\n")
                        break
            else:
                content_parts.append("(空文件)")
            
            full_text = "\n".join(content_parts)
            logger.info(f"CSV文件解析成功: {file_name}, 编码: {used_encoding}, 行数: {len(rows)}")
            return full_text
            
        except Exception as e:
            logger.error(f"CSV文件解析错误: {e}")
            return f"CSV文件解析失败: {str(e)}"
    
    async def _process_excel_file(self, file_path: Path, file_name: str) -> str:
        """处理Excel文件"""
        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)
            content_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"[工作表: {sheet_name}]")
                
                # 获取表格内容
                table_data = []
                for row in sheet.iter_rows(values_only=True):
                    # 过滤空行
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        table_data.append(row_text)
                
                if table_data:
                    content_parts.append("\n".join(table_data))
                else:
                    content_parts.append("(空表)")
                
                content_parts.append("")  # 添加空行分隔
            
            full_text = "\n".join(content_parts)
            logger.info(f"Excel文件解析成功: {file_name}, 工作表数: {len(workbook.sheetnames)}")
            return full_text
            
        except Exception as e:
            logger.error(f"Excel文件解析错误: {e}")
            return f"Excel文件解析失败: {str(e)}"
    
    async def _process_pptx_file(self, file_path: Path, file_name: str) -> str:
        """处理PowerPoint文件"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_count = len(prs.slides)
            
            content_parts = []
            content_parts.append(f"[PowerPoint文件]")
            content_parts.append(f"文件名: {file_name}")
            content_parts.append(f"幻灯片数量: {slides_count}")
            content_parts.append("")
            
            has_images = False
            has_tables = False
            has_notes = False
            
            # 收集待处理的图片
            images_to_process = []
            image_placeholders = {}
            image_count = 0
            
            # 遍历每一页幻灯片
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"=== 幻灯片 {slide_idx} ===")
                
                # 提取标题
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                    if title:
                        slide_content.append(f"标题: {title}")
                
                # 提取文本内容和图片
                text_content = []
                slide_images = []  # 该幻灯片的图片
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 避免重复添加标题
                        if shape != slide.shapes.title:
                            text_content.append(shape.text.strip())
                    
                    # 检测并提取图片
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True
                        try:
                            # 获取图片数据
                            image = shape.image
                            image_bytes = image.blob
                            
                            image_count += 1
                            placeholder = f"[[PPT_IMAGE_PLACEHOLDER_{image_count}]]"
                            
                            # 记录图片信息
                            images_to_process.append({
                                'index': image_count,
                                'slide_idx': slide_idx,
                                'placeholder': placeholder,
                                'data': image_bytes
                            })
                            
                            slide_images.append(placeholder)
                        except Exception as e:
                            logger.warning(f"无法提取幻灯片 {slide_idx} 的图片: {e}")
                    
                    # 检测表格
                    if hasattr(shape, 'table'):
                        has_tables = True
                        # 提取表格内容
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(" | ".join(row_data))
                        if table_data:
                            text_content.append("\n表格内容:")
                            text_content.extend(table_data)
                
                if text_content:
                    slide_content.append("内容:")
                    slide_content.extend(text_content)
                
                # 添加图片占位符
                if slide_images:
                    slide_content.extend(slide_images)
                
                # 提取演讲者备注
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    has_notes = True
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    slide_content.append(f"备注: {notes}")
                
                slide_content.append("")  # 添加空行分隔
                content_parts.extend(slide_content)
            
            # 如果有图片需要处理，并发处理所有图片
            if images_to_process and settings.VISION_API_KEY:
                from .multimodal_service import multimodal_service
                logger.info(f"PowerPoint文档包含 {len(images_to_process)} 张图片，开始并发处理...")
                
                # 创建并发任务
                async def process_single_image(image_info):
                    """处理单张图片的异步函数"""
                    index = image_info['index']
                    slide_idx = image_info['slide_idx']
                    placeholder = image_info['placeholder']
                    try:
                        result = await multimodal_service.extract_image_content(
                            image_info['data'],
                            options={
                                'prompt': '请详细描述这张图片的内容，如果是图表请提取其中的数据和信息，如果有文字请准确识别。'
                            }
                        )
                        
                        if result['success']:
                            logger.info(f"成功提取幻灯片 {slide_idx} 图片 {index} 内容")
                            return {
                                'placeholder': placeholder,
                                'content': f"\n[图片 {index}]\n{result['content']}\n",
                                'success': True
                            }
                        else:
                            logger.warning(f"图片 {index} 提取失败: {result.get('error', '未知错误')}")
                            return {
                                'placeholder': placeholder,
                                'content': f"\n[图片 {index}：无法识别 - {result.get('error', '未知错误')}]\n",
                                'success': False
                            }
                    except Exception as e:
                        logger.error(f"处理图片 {index} 时出错: {e}")
                        return {
                            'placeholder': placeholder,
                            'content': f"\n[图片 {index}：处理失败 - {str(e)}]\n",
                            'success': False
                        }
                
                # 限制并发数，避免API限流
                import asyncio
                MAX_CONCURRENT = settings.IMAGE_PROCESS_MAX_CONCURRENT
                semaphore = asyncio.Semaphore(MAX_CONCURRENT)
                
                async def process_with_semaphore(image_info):
                    async with semaphore:
                        try:
                            return await asyncio.wait_for(
                                process_single_image(image_info), 
                                timeout=settings.IMAGE_PROCESS_TIMEOUT
                            )
                        except asyncio.TimeoutError:
                            logger.error(f"处理图片 {image_info['index']} 超时")
                            return {
                                'placeholder': image_info['placeholder'],
                                'content': f"\n[图片 {image_info['index']}：处理超时]\n",
                                'success': False
                            }
                
                # 并发处理所有图片
                tasks = [process_with_semaphore(img) for img in images_to_process]
                results = await asyncio.gather(*tasks, return_exceptions=True)
                
                # 建立占位符到内容的映射
                placeholder_map = {}
                success_count = 0
                for result in results:
                    if isinstance(result, Exception):
                        logger.error(f"图片处理任务异常: {result}")
                    elif isinstance(result, dict):
                        placeholder_map[result['placeholder']] = result['content']
                        if result.get('success'):
                            success_count += 1
                
                logger.info(f"图片并发处理完成，成功处理 {success_count} 张图片")
                
                # 替换占位符为实际内容
                for i, part in enumerate(content_parts):
                    if part in placeholder_map:
                        content_parts[i] = placeholder_map[part]
            elif images_to_process and not settings.VISION_API_KEY:
                # 如果有图片但没有配置视觉模型，替换占位符为提示信息
                for i, part in enumerate(content_parts):
                    if part.startswith("[[PPT_IMAGE_PLACEHOLDER_"):
                        image_num = part.split('_')[3].rstrip(']]')
                        content_parts[i] = f"\n[图片 {image_num}：未配置视觉模型]\n"
            
            # 添加文件特性摘要
            features = []
            if has_images:
                features.append("包含图片")
            if has_tables:
                features.append("包含表格")
            if has_notes:
                features.append("包含演讲者备注")
            
            if features:
                # 在文件基本信息后面插入特性摘要
                for i, part in enumerate(content_parts):
                    if part == "":  # 找到第一个空行
                        content_parts.insert(i, f"文件特性: {', '.join(features)}")
                        break
            
            full_text = "\n".join(content_parts)
            logger.info(f"PowerPoint文件解析成功: {file_name}, 幻灯片数: {slides_count}")
            return full_text
            
        except ImportError:
            logger.warning(f"缺少python-pptx库，无法解析PowerPoint文件: {file_name}")
            return f"[PowerPoint文件: {file_name}]\n\n需要安装 python-pptx 库才能解析PowerPoint文件。请运行: pip install python-pptx"
        except Exception as e:
            logger.error(f"PowerPoint文件解析错误: {e}")
            return f"PowerPoint文件解析失败: {str(e)}"
    
    async def _process_image_file(self, file_path: Path, file_name: str) -> str:
        """处理图片文件"""
        logger.info(f"开始处理图片文件: {file_name}")
        
        # 检查是否配置了视觉模型
        if settings.VISION_API_KEY:
            try:
                from .multimodal_service import multimodal_service
                
                # 读取图片文件内容
                with open(file_path, 'rb') as f:
                    file_content = f.read()
                
                # 调用视觉模型分析图片
                result = await multimodal_service.extract_image_content(
                    image_data=file_content,
                    options={
                        'prompt': '请详细描述这张图片的内容，包括其中的文字、图表、数据等所有信息。'
                    }
                )
                
                if result['success']:
                    logger.info(f"图片分析成功: {file_name}")
                    return f"[图片文件: {file_name}]\n\n{result['content']}"
                else:
                    logger.warning(f"图片分析失败: {file_name}, 错误: {result.get('error')}")
                    return f"[图片文件: {file_name}]\n\n图片分析失败: {result.get('error', '未知错误')}"
            except Exception as e:
                logger.error(f"图片处理异常: {e}", exc_info=True)
                return f"[图片文件: {file_name}]\n\n图片处理出错: {str(e)}"
        else:
            return f"[图片文件: {file_name}]\n\n未配置视觉模型，无法分析图片内容。请配置 VISION_API_KEY。"
    
    def _handle_unsupported_file(self, file_ext: str, file_name: str) -> str:
        """处理不支持的文件格式"""
        unsupported_msg = f"[{file_ext} 文件: {file_name}]\n\n"
        
        # 检查是否是因为缺少依赖
        if file_ext == '.xlsx' and not HAS_OPENPYXL:
            unsupported_msg += "Excel 文件解析功能未启用。请联系管理员安装 openpyxl 库。"
            logger.warning(f"用户尝试上传 Excel 文件，但未安装 openpyxl 库: {file_name}")
        elif file_ext == '.pdf' and not HAS_PDF:
            unsupported_msg += "PDF 文件解析功能未启用。请联系管理员安装 pypdf 库。"
            logger.warning(f"用户尝试上传 PDF 文件，但未安装 pypdf 库: {file_name}")
        elif file_ext == '.docx' and not HAS_DOCX:
            unsupported_msg += "Word 文档解析功能未启用。请联系管理员安装 python-docx 库。"
            logger.warning(f"用户尝试上传 Word 文档，但未安装 python-docx 库: {file_name}")
        else:
            unsupported_msg += "暂不支持此格式的文档解析。"
        
        return unsupported_msg
    
    async def _dispatch_file_processing(self, file_path: Path, file_ext: str, file_name: str, file_id: str) -> str:
        """文件处理分发器 - 根据文件类型调用相应的处理器"""
        # 获取对应的处理器
        processor = self._file_processors.get(file_ext.lower())
        
        if processor is None:
            # 检查是否是因为缺少依赖导致处理器为None
            if file_ext.lower() in ['.pdf', '.docx', '.xlsx']:
                return self._handle_unsupported_file(file_ext, file_name)
            else:
                return self._handle_unsupported_file(file_ext, file_name)
        
        # 特殊处理.doc文件（返回三元组）
        if file_ext == '.doc':
            result = await processor(file_path, file_name, file_id)
            if isinstance(result, tuple):
                full_text, new_file_path, new_file_ext = result
                # 如果.doc转换成功，继续处理.docx
                if new_file_ext == '.docx' and full_text == "":
                    docx_processor = self._file_processors.get('.docx')
                    if docx_processor:
                        return await docx_processor(new_file_path, file_name)
                return full_text
            return result
        
        # 普通文件处理
        return await processor(file_path, file_name)
    
    def _create_document_chunks(self, text: str) -> List[str]:
        """将文档文本分块"""
        # 分块处理 - 优化分块策略
        # 对于大文档，使用更大的分块尺寸以减少分块数量
        base_chunk_size = 1000
        if len(text) > 100000:  # 超过10万字符的大文档
            chunk_size = 2000  # 使用更大的分块
            chunk_overlap = 400
        elif len(text) > 50000:  # 5万-10万字符的中等文档
            chunk_size = 1500
            chunk_overlap = 300
        else:  # 小文档
            chunk_size = base_chunk_size
            chunk_overlap = 200
        
        chunks = []
        
        if len(text) <= chunk_size:
            chunks = [text]
        else:
            for i in range(0, len(text), chunk_size - chunk_overlap):
                chunk = text[i:i + chunk_size]
                chunks.append(chunk)
        
        return chunks
    
    async def _finalize_document_processing(self, db: AsyncSession, file_id: str, full_text: str, file_ext: str) -> int:
        """完成文档处理，分块并更新数据库"""
        # 创建分块
        chunks = self._create_document_chunks(full_text)
        
        # 提取元数据
        metadata = {
            "char_count": len(full_text),
            "chunk_count": len(chunks),
            "file_type": file_ext
        }
        
        # 更新数据库
        async with db.begin():
            # 重新查询文档记录
            result = await db.execute(
                select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
            )
            doc_upload = result.scalar_one_or_none()
            if doc_upload:
                doc_upload.process_status = ProcessStatus.READY
                # MEDIUMTEXT 支持 16MB，约 1600万字符，这里限制为 500万字符以保留余量
                doc_upload.doc_content = full_text[:5000000] if len(full_text) > 5000000 else full_text
                
                # 限制分块数量，避免 JSON 过大
                # 每个分块 1000 字符，存储前 1000 个分块（约 100万字符的内容）
                max_chunks = 1000
                chunks_to_store = chunks[:max_chunks]
                doc_upload.doc_chunks = json.dumps([{"id": i, "content": chunk} for i, chunk in enumerate(chunks_to_store)])
                
                doc_upload.doc_metadata = json.dumps(metadata)
                doc_upload.process_end_time = now_shanghai()
                await db.flush()
        
        return len(full_text)

    async def _process_document_async(self, file_id: str) -> None:
        """
        异步处理文档的包装方法
        """
        try:
            # 使用独立的数据库会话
            from src.shared.db.config import get_async_db_context
            async with get_async_db_context() as db:
                await self.process_document(db, file_id)
        except Exception as e:
            logger.error(f"异步处理文档失败: {file_id}, 错误: {str(e)}")

    async def get_file_info(self, db: AsyncSession, file_id: str, user_name: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """
        获取文件信息（用于下载）
        
        Args:
            db: 数据库会话
            file_id: 文件ID
            user_name: 用户名（用于权限检查）
            
        Returns:
            文件信息字典，包含 file_path 和 file_name
        """
        # 查询文档
        query = select(AgentDocumentUpload).where(AgentDocumentUpload.file_id == file_id)
        
        # 如果提供了用户名，检查所有权
        if user_name:
            query = query.where(AgentDocumentUpload.create_by == user_name)
        
        result = await db.execute(query)
        doc_upload = result.scalar_one_or_none()
        
        if not doc_upload:
            return None
        
        # 构建文件路径，需要加上文件扩展名
        file_path = os.path.join(UPLOAD_DIR, f"{doc_upload.file_id}{doc_upload.file_type}")
        
        return {
            "file_id": doc_upload.file_id,
            "file_name": doc_upload.file_name,
            "file_path": file_path,
            "file_type": doc_upload.file_type,
            "file_size": doc_upload.file_size,
            "upload_time": doc_upload.upload_time.isoformat() if doc_upload.upload_time else None
        }


# 全局文档服务实例
document_service = DocumentService()